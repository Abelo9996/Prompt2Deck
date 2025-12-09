"""
Slide Builder Module
Assembles slides into a PowerPoint presentation using python-pptx.
"""

import os
from typing import List, Optional
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from models import SlideData


class SlideBuilder:
    """
    Builds PowerPoint presentations from slide data.
    Handles formatting, layout, and export.
    """
    
    def __init__(self):
        """Initialize the slide builder."""
        self.themes = {
            "professional": {
                "title_color": RGBColor(31, 56, 100),  # Dark blue
                "text_color": RGBColor(64, 64, 64),    # Dark gray
                "background": RGBColor(255, 255, 255),  # White
                "accent": RGBColor(74, 144, 226)        # Blue
            },
            "modern": {
                "title_color": RGBColor(41, 128, 185),  # Bright blue
                "text_color": RGBColor(44, 62, 80),     # Dark slate
                "background": RGBColor(236, 240, 241),  # Light gray
                "accent": RGBColor(230, 126, 34)        # Orange
            },
            "minimal": {
                "title_color": RGBColor(0, 0, 0),       # Black
                "text_color": RGBColor(80, 80, 80),     # Gray
                "background": RGBColor(255, 255, 255),  # White
                "accent": RGBColor(100, 100, 100)       # Gray
            }
        }
    
    async def build_deck(
        self,
        slides: List[SlideData],
        output_dir: str,
        theme: str = "professional"
    ) -> str:
        """
        Build a PowerPoint presentation from slide data.
        
        Args:
            slides: List of slides with content
            output_dir: Directory to save the presentation
            theme: Visual theme to apply
            
        Returns:
            Path to the generated PPTX file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        theme_colors = self.themes.get(theme, self.themes["professional"])
        
        # Add title slide
        self._add_title_slide(prs, slides[0], theme_colors)
        
        # Add content slides
        for slide_data in slides[1:]:
            self._add_content_slide(prs, slide_data, theme_colors)
        
        # Generate filename and save
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        filepath = os.path.join(output_dir, filename)
        
        prs.save(filepath)
        return filepath
    
    def _add_title_slide(
        self,
        prs: Presentation,
        slide_data: SlideData,
        theme_colors: dict
    ):
        """Add a title slide to the presentation."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = slide_data.title
        
        # Format title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_colors["title_color"]
        
        # Add subtitle if bullets exist
        if slide_data.bullets:
            subtitle_text = " | ".join(slide_data.bullets[:3])
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(8), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = theme_colors["text_color"]
    
    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: SlideData,
        theme_colors: dict
    ):
        """Add a content slide to the presentation."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.title
        
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_colors["title_color"]
        
        # Add bullets
        if slide_data.bullets:
            bullet_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5)
            )
            text_frame = bullet_box.text_frame
            text_frame.word_wrap = True
            
            for i, bullet in enumerate(slide_data.bullets):
                if i > 0:
                    text_frame.add_paragraph()
                
                para = text_frame.paragraphs[i]
                para.text = bullet
                para.level = 0
                para.font.size = Pt(20)
                para.font.color.rgb = theme_colors["text_color"]
                para.space_before = Pt(12)
        
        # Add speaker notes if present
        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.speaker_notes
    
    async def export_to_pdf(self, pptx_path: str) -> Optional[str]:
        """
        Export presentation to PDF.
        
        Note: This requires LibreOffice or PowerPoint to be installed.
        Alternatively, use a PDF conversion library.
        
        Args:
            pptx_path: Path to PPTX file
            
        Returns:
            Path to PDF file or None if export failed
        """
        try:
            # Try using LibreOffice for conversion
            pdf_path = pptx_path.replace('.pptx', '.pdf')
            
            # This is a placeholder - actual implementation would need:
            # - LibreOffice headless mode
            # - Or a Python PDF library
            # - Or a cloud conversion service
            
            import subprocess
            result = subprocess.run([
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(pptx_path),
                pptx_path
            ], capture_output=True, timeout=30)
            
            if result.returncode == 0 and os.path.exists(pdf_path):
                return pdf_path
            else:
                print("PDF export failed: LibreOffice not available or conversion error")
                return None
                
        except Exception as e:
            print(f"PDF export error: {e}")
            return None
